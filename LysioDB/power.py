import os
import pptx
import re
import numpy as np
import polars as pl
from typing import Dict, Tuple, List
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.chart import Chart
from pptx.chart.data import ChartData, XyChartData


class Power:
    def __init__(self, database):
        """
        Initializes the Builder with necessary data.

        Args:
            database (database): An instance of the database class providing access to data.
        """
        self.database = database
        self.matrix_counter = 0
        self.min_corr = None
        self.max_corr = None
        self.min_index = None
        self.max_index = None

    def update_pptx(
        self,
        category,
        length: int = 0,
        template_path: str = "template.ppt",
        output_dir: str = "output_ppts",
    ):
        """Generate a PowerPoint for each unique category."""

        print("\n--- Start generating Powerpoints ---")
        os.makedirs(output_dir, exist_ok=True)

        prs = pptx.Presentation(template_path)

        area_mapping = {
            f"area_{i}": area
            for i, area in enumerate(self.database.config.area_map.keys())
        }
        year_mapping = {
            f"year_{i}": year for i, year in enumerate(self.database.config.year_map)
        }
        self.matrix_counter = 0
        self.min_corr = None
        self.max_corr = None
        self.min_index = None
        self.max_index = None

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._parse_text_frame(
                        shape.text_frame, area_mapping, year_mapping, category
                    )

                elif shape.has_table:
                    self._update_table(
                        shape.table, area_mapping, year_mapping, category
                    )
                    if length > 0:
                        self._process_table(shape.table, length)

                elif shape.has_chart:
                    self._update_chart(shape.chart, category)

        output_path = f"{output_dir}/powerpoint_{category}.pptx"
        prs.save(output_path)

        print(f"Generated PowerPoints in {output_dir}")

    def _process_table(self, table, length: int) -> None:
        """Process all cells in a table to highlight min/max values."""
        table_values = self._collect_table_values(table, length)

        for row_idx, row_data in table_values.items():
            min_vals, max_vals = row_data
            for col_idx, cell in enumerate(table.rows[row_idx].cells):
                self._format_cell_if_extreme(cell, min_vals, max_vals, length)

    def _collect_table_values(
        self, table, length: int
    ) -> Dict[int, Tuple[List[float], List[float]]]:
        """Collect all percentage values and calculate min/max for each row."""
        table_values = {}

        for row_idx, row in enumerate(table.rows):
            values = []
            for cell in row.cells:
                if match := re.match(r"(\d+)%", cell.text):
                    values.append(float(match.group(1)))

            if len(values) > 1:
                sorted_values = sorted(values)
                table_values[row_idx] = (
                    sorted_values[:length],
                    sorted_values[-length:],
                )
            else:
                table_values[row_idx] = ([], [])

        return table_values

    def _format_cell_if_extreme(
        self, cell, min_vals: List[float], max_vals: List[float], length: int
    ) -> None:
        """Format cell if it contains a min or max value."""
        if match := re.match(r"(\d+)%", cell.text):
            numeric_value = float(match.group(1))

            if numeric_value in min_vals:
                self._format_cell(cell, f"• {cell.text}", RGBColor(240, 120, 100))
            elif numeric_value in max_vals:
                self._format_cell(cell, f"• {cell.text}", RGBColor(125, 186, 116))

    def _format_cell(self, cell, text: str, color: RGBColor) -> None:
        """Apply consistent formatting to a cell."""
        cell.text = text
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]

        for run in paragraph.runs:
            run.text = ""

        run = paragraph.add_run()
        run.text = text
        paragraph.alignment = PP_ALIGN.CENTER

        font = run.font
        font.size = Pt(12)
        font.bold = True
        font.color.rgb = color

    def _parse_text_frame(self, text_frame, area_mapping, year_mapping, category):
        """Replace simple placeholders in text frames."""

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                placeholders = self._extract_placeholders(text)

                for placeholder in placeholders:
                    cat_match = re.match(r"category", placeholder)
                    freq_match = re.match(r"frequency", placeholder)
                    area_year_match = re.match(r"area_(\d+):year_(\d+)", placeholder)

                    if area_year_match:
                        year = area_year_match.group(0).split(":")[-1]
                        area = area_year_match.group(0).split(":")[0]
                        year_value = year_mapping.get(year)
                        area_name = area_mapping.get(area)
                        category_label = f"{year_value}:{category}"
                        if category_label in self.database.index_df["Category"]:
                            try:
                                result_value = (
                                    self.database.index_df.filter(
                                        pl.col("Category") == category_label
                                    )
                                    .select(pl.col(area_name).round(2))
                                    .item(0, 0)
                                )
                                result_value = self._format_number_swedish(result_value)

                                text = str(result_value)

                            except (IndexError, KeyError) as e:
                                print(
                                    f"Warning: Could not find data for {placeholder} in category '{category}' ({e})"
                                )

                    if cat_match:
                        text = category

                    if freq_match:
                        text = "test"
                run.text = text

    def _update_table(self, table, area_mapping, year_mapping, category):
        """Update data in existing tables based on placeholders."""
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text
                placeholders = self._extract_placeholders(text)
                for placeholder in placeholders:
                    question_pattern = "|".join(
                        rf"{prefix}(.*)"
                        for prefix in self.database.config.QUESTION_PREFIXES
                    )
                    area_year_match = re.match(r"area_(\d+):year_(\d+)", placeholder)
                    area_match = re.match(r"area_(\d+)", placeholder)
                    year_match = re.match(r"year_(\d+)", placeholder)
                    question_match = re.search(
                        question_pattern,
                        placeholder,
                    )
                    nan_percentage_match = re.search("nan:(.*)", placeholder)
                    count_match = re.search(r"count:(.*)", placeholder)

                    if area_match:
                        area = area_match.group(0)
                        area_name = area_mapping.get(area)
                        self._update_cell_text(cell, str(area_name))

                    if year_match:
                        year_index = int(year_match.group(1))
                        year = year_match.group(0)
                        if 0 <= year_index < len(self.database.config.year_map):
                            year_value = year_mapping.get(year)
                            self._update_cell_text(cell, str(year_value))

                    if question_match:
                        year_match = re.search(":", question_match.group(0))
                        if year_match:
                            question = question_match.group(0).split(":")[0]
                            year = question_match.group(0).split(":")[-1]
                            year_value = year_mapping.get(year)
                            category_label = f"{year_value}:{category}"
                            if category_label in self.database.index_df["Category"]:
                                try:
                                    result_value = (
                                        self.database.index_df.filter(
                                            pl.col("Category") == category_label
                                        )
                                        .select(pl.col(question).round(2))
                                        .item(0, 0)
                                    )
                                    result_value = self._format_number_swedish(
                                        result_value
                                    )
                                    if result_value == "nan":
                                        result_value = "-"
                                    self._update_cell_text(cell, str(result_value))
                                except (IndexError, KeyError) as e:
                                    print(
                                        f"Warning: Could not find data for {placeholder} in category '{category}' ({e})"
                                    )
                        else:
                            question_part = question_match.group(0)
                            base_question = self._get_base_question(question_part)
                            question_label = self._get_column_label_for_sub_question(
                                base_question, question_part
                            )

                            self._update_cell_text(cell, str(question_label))

                    if area_year_match:
                        year = area_year_match.group(0).split(":")[-1]
                        area = area_year_match.group(0).split(":")[0]
                        year_value = year_mapping.get(year)
                        area_name = area_mapping.get(area)
                        category_label = f"{year_value}:{category}"
                        if category_label in self.database.index_df["Category"]:
                            try:
                                result_value = (
                                    self.database.index_df.filter(
                                        pl.col("Category") == category_label
                                    )
                                    .select(pl.col(area_name).round(2))
                                    .item(0, 0)
                                )
                                result_value = self._format_number_swedish(result_value)
                                if result_value == "nan":
                                    result_value = "-"

                                self._update_cell_text(cell, str(result_value))
                            except (IndexError, KeyError) as e:
                                print(
                                    f"Warning: Could not find data for {placeholder} in category '{category}' ({e})"
                                )

                    if nan_percentage_match:
                        question_part = nan_percentage_match.group(0).split(":")[-1]
                        year = year_mapping.get("year_0")
                        base_question = self._get_base_question(question_part)

                        try:
                            nan = (
                                self.database.percentage_df.filter(
                                    (pl.col("question") == question_part)
                                    & (pl.col("metric_type") == "percentage")
                                    & (pl.col("answer_value") == "nan")
                                )
                                .select(pl.col(f"{year}:{category}"))
                                .item(0, 0)
                            )

                            if nan is not None:
                                nan_percentage_value = nan * 100
                                self._update_cell_text(
                                    cell, f"{int(nan_percentage_value)}%"
                                )
                            else:
                                self._update_cell_text(cell, "0%")

                        except (KeyError, IndexError) as e:
                            print(
                                f"Warning: Could not find nan_percentage for {placeholder} in category '{category}' ({e})"
                            )
                            self._update_cell_text(cell, "0%")

                    if count_match:
                        question_part = count_match.group(0).split(":")[-1]
                        year = year_mapping.get("year_0")
                        base_question = self._get_base_question(question_part)
                        try:
                            count_value = int(
                                self.database.percentage_df.filter(
                                    (pl.col("question") == question_part)
                                    & (pl.col("metric_type") == "count")
                                    & (pl.col("answer_value") != "nan")
                                    & (pl.col("answer_value") != "total")
                                )
                                .select(pl.col(f"{year}:{category}"))
                                .sum()
                                .item(0, 0)
                            )

                            self._update_cell_text(cell, str(count_value))
                        except (KeyError, IndexError) as e:
                            print(
                                f"Warning: Could not find count for {placeholder} in category '{category}' ({e})"
                            )
                            self._update_cell_text(cell, "N/A")

    def _update_cell_text(self, cell, new_text):
        """Updates the text of a cell, preserving format of the first run."""
        text_frame = cell.text_frame
        if len(text_frame.paragraphs) > 0 and len(text_frame.paragraphs[0].runs) > 0:
            text_frame.paragraphs[0].runs[0].text = str(new_text)
        else:
            cell.text = str(new_text)

    def _update_chart(self, chart, category):
        chart_data = ChartData()
        year = self.database.config.year_map[0]
        category = f"{year}:{category}"

        if isinstance(chart, Chart):
            if chart.chart_type in [
                XL_CHART_TYPE.PIE,
                XL_CHART_TYPE.LINE,
                XL_CHART_TYPE.LINE_MARKERS,
                XL_CHART_TYPE.PIE_EXPLODED,
            ]:
                template_var_names = chart.plots[0].categories.flattened_labels
                var_names = tuple(label[0] for label in template_var_names)
                template_val_name_suffixes = tuple(
                    series.name for series in chart.series
                )

                var_labels = []
                all_value_labels = []
                all_series_values = {}
                for suffix in template_val_name_suffixes:
                    question_part = suffix
                    base_question = self._get_base_question(question_part)
                    var_labels.append(
                        self._get_column_label_for_sub_question(
                            base_question, question_part
                        )
                    )
                    series_values = []
                    for var in var_names:
                        if var != "":
                            try:
                                var_num = float(var)
                                value = (
                                    self.database.percentage_df.filter(
                                        (pl.col("question") == question_part)
                                        & (pl.col("answer_value") == str(var_num))
                                        & (pl.col("metric_type") == "percentage")
                                    )
                                    .select(category)
                                    .item(0, 0)
                                )
                                if not value:
                                    series_values.append("nan")
                                else:
                                    series_values.append(value)
                                if not all_value_labels:
                                    all_value_labels.append(
                                        self.database.question_df.filter(
                                            pl.col("question") == question_part
                                        )
                                        .select(pl.col("value_labels_info"))
                                        .item(0, 0)
                                    )
                            except KeyError as e:
                                print(f"KeyError: {e}")
                            except ValueError as e:
                                print(f"ValueError: {e}")

                    all_series_values[suffix] = series_values

                chart_data.categories = all_value_labels[0].values()
                for series_label, series_data in all_series_values.items():
                    chart_data.add_series(series_label, series_data)

                chart.replace_data(chart_data)

            if chart.chart_type in [
                XL_CHART_TYPE.BAR_STACKED_100,
                XL_CHART_TYPE.BAR_CLUSTERED,
                XL_CHART_TYPE.COLUMN_CLUSTERED,
            ]:
                template_var_names = chart.plots[0].categories.flattened_labels
                var_names = tuple(label[0] for label in template_var_names)
                template_val_name_suffixes = tuple(
                    int(series.name) for series in chart.series
                )

                var_labels = []
                all_value_labels = []
                all_series_values = []
                for var in var_names:
                    question_part = var
                    base_question = self._get_base_question(question_part)
                    var_labels.append(
                        self._get_column_label_for_sub_question(
                            base_question, question_part
                        )
                    )
                    series_values = []
                    for suffix in template_val_name_suffixes:
                        try:
                            var_num = float(suffix)
                            value = (
                                self.database.percentage_df.filter(
                                    (pl.col("question") == question_part)
                                    & (pl.col("answer_value") == str(var_num))
                                    & (pl.col("metric_type") == "percentage")
                                )
                                .select(category)
                                .item(0, 0)
                            )
                            if not value:
                                series_values.append("nan")
                            else:
                                series_values.append(value)
                            if not all_value_labels:
                                all_value_labels.append(
                                    self.database.question_df.filter(
                                        pl.col("question") == question_part
                                    )
                                    .select(pl.col("value_labels_info"))
                                    .item(0, 0)
                                )
                        except KeyError as e:
                            print(f"KeyError: {e}")
                        except ValueError as e:
                            print(f"ValueError: {e}")

                    all_series_values.append(series_values)

                chart_data.categories = var_labels
                for i, suffix in enumerate(template_val_name_suffixes):
                    try:
                        series_label = all_value_labels[0].get(str(float(i) + 1))

                        series_data = [row[i] for row in all_series_values]

                        chart_data.add_series(series_label, series_data)

                    except (KeyError, ValueError, IndexError) as e:
                        print(f"Error adding series: {e}")

                chart.replace_data(chart_data)

            elif chart.chart_type == XL_CHART_TYPE.XY_SCATTER:
                chart_data = XyChartData()
                all_corr_values = []
                all_index_values = []
                if self.max_index is None:
                    self._get_min_max(category)

                chart.value_axis.minimum_scale = max(self.min_index - 0.2, 0)
                chart.value_axis.maximum_scale = min(self.max_index + 0.2, 5)
                chart.category_axis.minimum_scale = max(self.min_corr - 0.02, 0)
                chart.category_axis.maximum_scale = min(self.max_corr + 0.02, 1)
                chart.category_axis.visible = False
                chart.value_axis.visible = False

                for i, question in self.database.matrix[self.matrix_counter].items():
                    index = (
                        self.database.index_df.filter(pl.col("Category") == category)
                        .select(question)
                        .item(0, 0)
                    )
                    corr = (
                        self.database.correlate_df.filter(
                            pl.col("Question") == question
                        )
                        .select("Correlation")
                        .item(0, 0)
                    )
                    if corr is None or corr <= 0:
                        corr = 0
                    if index is None:
                        index = 0

                    all_corr_values.append(corr)
                    all_index_values.append(index)

                    matrix_series = chart_data.add_series(str(i))
                    matrix_series.add_data_point(corr, index)

                self.matrix_counter += 1
                chart.replace_data(chart_data)

            else:
                print(f"Warning: Chart type {chart.chart_type} is not supported.")

        else:
            print(
                f"Warning: Shape's chart attribute for category '{category}' is not a Chart object, but a {type(chart)}."
            )

    def _extract_placeholders(self, text):
        """Extract all placeholders from the text using regex."""
        brackets = re.findall(r"{(.*?)}", text)
        category = re.findall(r"category:category", text)
        if brackets:
            return brackets
        elif category:
            return category
        return [""]

    def _get_column_label_for_sub_question(self, base_question, column):
        """
        Retrieves the specific column label for a sub-question in a multi-response question.

        Args:
            base_question (str): The base question identifier (e.g., 'Q19').
            column (str): The specific column name (e.g., 'Q19C1').

        Returns:
            str or None: The column label for the given sub-question, or None if not found.
        """
        if base_question in self.database.question_df["base_question"]:
            try:
                return self.database.question_df.filter(pl.col("question") == column)[
                    "question_label"
                ][0]
            except ValueError:
                return None
        return None

    def _format_number_swedish(self, number):
        """Formats a number using Swedish decimal separator (comma)."""
        if isinstance(number, (int, float)):
            return str(number).replace(".", ",")
        return str(number)

    def _get_base_question(self, question):
        if question in self.database.question_df["question"]:
            base_question = (
                self.database.question_df.filter(pl.col("question") == question)
                .select(pl.col("base_question"))
                .item(0, 0)
            )
        return str(base_question)

    def _get_min_max(self, category):
        area_map = self.database.config.area_map.copy()
        first_key = next(iter(area_map))
        del area_map[first_key]

        all_index = []
        all_corr = []
        for area, questions in area_map.items():
            for question in questions:
                index = (
                    self.database.index_df.filter(pl.col("Category") == category)
                    .select(area)
                    .item(0, 0)
                )

                corr = (
                    self.database.correlate_df.filter(pl.col("Question") == question)
                    .select("Correlation")
                    .item(0, 0)
                )
                if (corr is None) or (index is None):
                    continue
                else:
                    all_index.append(index)
                    all_corr.append(corr)

        if not all_index:
            all_index.append(0)
            all_index.append(5)

        if not all_corr:
            all_corr.append(0)
            all_corr.append(1)

        self.max_corr = max(all_corr)
        self.min_corr = min(all_corr)
        self.min_index = min(all_index)
        self.max_index = max(all_index)

    def _process_table(self, table, length: int) -> None:
        """Process all cells in a table to highlight min/max values."""
        table_values = self._collect_table_values(table, length)

        for row_idx, row_data in table_values.items():
            min_vals, max_vals = row_data
            for col_idx, cell in enumerate(table.rows[row_idx].cells):
                self._format_cell_if_extreme(cell, min_vals, max_vals, length)

    def _collect_table_values(
        self, table, length: int
    ) -> Dict[int, Tuple[List[float], List[float]]]:
        """Collect all percentage values and calculate min/max for each row."""
        table_values = {}

        for row_idx, row in enumerate(table.rows):
            values = []
            for cell in row.cells:
                if match := re.match(r"(\d+)%", cell.text):
                    values.append(float(match.group(1)))

            if len(values) > 1:
                sorted_values = sorted(values)
                table_values[row_idx] = (
                    sorted_values[:length],
                    sorted_values[-length:],
                )
            else:
                table_values[row_idx] = ([], [])

        return table_values

    def _format_cell_if_extreme(
        self, cell, min_vals: List[float], max_vals: List[float], length: int
    ) -> None:
        """Format cell if it contains a min or max value."""
        if match := re.match(r"(\d+)%", cell.text):
            numeric_value = float(match.group(1))

            if numeric_value in min_vals:
                self._format_cell(cell, f"• {cell.text}", RGBColor(240, 120, 100))
            elif numeric_value in max_vals:
                self._format_cell(cell, f"• {cell.text}", RGBColor(125, 186, 116))

    def _format_cell(self, cell, text: str, color: RGBColor) -> None:
        """Apply consistent formatting to a cell."""
        cell.text = text
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]

        for run in paragraph.runs:
            run.text = ""

        run = paragraph.add_run()
        run.text = text
        paragraph.alignment = PP_ALIGN.CENTER

        font = run.font
        font.size = Pt(12)
        font.bold = True
        font.color.rgb = color
